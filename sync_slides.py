"""
Sync Slides: Reads updated Financial Model XLSX → updates Pitch Deck PPTX.
This is a service-compatible version of the existing sync_model_to_slides.py.
"""

from pathlib import Path
from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData


def fmt_k(val):
    if val is None or val == 0:
        return "$0"
    if abs(val) < 1_000_000:
        return f"${val / 1000:.0f}K"
    return f"${val / 1_000_000:.2f}M"

def fmt_comma(val):
    if val is None:
        return "0"
    return f"{int(val):,}"

def fmt_pct(val):
    if val is None:
        return "0%"
    return f"{int(val * 100)}%"


def replace_text_in_shape(shape, old_text, new_text):
    if not shape.has_text_frame:
        return False
    found = False
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)
                found = True
        if found:
            continue
        full = ''.join(r.text for r in paragraph.runs)
        if old_text in full:
            new_full = full.replace(old_text, new_text)
            if paragraph.runs:
                paragraph.runs[0].text = new_full
                for r in paragraph.runs[1:]:
                    r.text = ''
                found = True
    return found


def extract_excel_data(excel_path):
    wb = load_workbook(excel_path, data_only=True)
    data = {}

    ws_a = wb["Assumptions"]
    data['pricing_light'] = ws_a['B7'].value or 100
    data['pricing_growth'] = ws_a['B8'].value or 180
    data['pricing_pro'] = ws_a['B9'].value or 350
    data['pricing_enterprise'] = ws_a['B10'].value or 999
    data['churn_rate'] = ws_a['B27'].value or 0.03

    prefund_totals = [ws_a[f'{c}17'].value or 0 for c in 'BCDE']
    data['prefund_quarterly'] = prefund_totals

    postfund_totals = [ws_a[f'{c}24'].value or 0 for c in 'FGHI']
    data['postfund_annual'] = postfund_totals

    ws_pl = wb["P&L"]
    for fy, col in [('fy26', 'F'), ('fy27', 'G'), ('fy28', 'H'), ('fy29', 'I')]:
        data[f'revenue_{fy}'] = ws_pl[f'{col}13'].value or 0
        data[f'costs_{fy}'] = ws_pl[f'{col}24'].value or 0
        data[f'net_income_{fy}'] = ws_pl[f'{col}26'].value or 0

    data['closing_cash_fy26'] = ws_pl['F31'].value or 0
    data['closing_cash_fy29'] = ws_pl['I31'].value or 0

    ws_w = wb["Customer Waterfall"]
    for fy, col in [('fy26', 'F'), ('fy27', 'G'), ('fy28', 'H'), ('fy29', 'I')]:
        data[f'light_customers_{fy}'] = ws_w[f'{col}9'].value or 0
        data[f'growth_customers_{fy}'] = ws_w[f'{col}15'].value or 0
        data[f'pro_customers_{fy}'] = ws_w[f'{col}21'].value or 0
        data[f'enterprise_customers_{fy}'] = ws_w[f'{col}27'].value or 0
        data[f'total_customers_{fy}'] = ws_w[f'{col}30'].value or 0

    ws_g = wb["Geographic Split"]
    for fy, col in [('fy26', 'B'), ('fy27', 'C'), ('fy28', 'D'), ('fy29', 'E')]:
        data[f'egypt_{fy}'] = ws_g[f'{col}10'].value or 0
        data[f'gcc_{fy}'] = ws_g[f'{col}17'].value or 0
        data[f'nafrica_{fy}'] = ws_g[f'{col}24'].value or 0
        data[f'global_total_{fy}'] = ws_g[f'{col}27'].value or 0

    data['light_by_fy'] = [ws_g[f'{c}6'].value or 0 for c in 'BCDE']
    data['growth_by_fy'] = [ws_g[f'{c}7'].value or 0 for c in 'BCDE']
    data['pro_by_fy'] = [ws_g[f'{c}8'].value or 0 for c in 'BCDE']
    data['enterprise_by_fy'] = [ws_g[f'{c}9'].value or 0 for c in 'BCDE']

    if data['total_customers_fy29'] > 0:
        blended_arpu = (data['revenue_fy29'] / 12) / data['total_customers_fy29']
    else:
        blended_arpu = 0

    if data['churn_rate'] > 0:
        months_to_churn = 1 / data['churn_rate']
        blended_ltv = blended_arpu * months_to_churn
    else:
        months_to_churn = 0
        blended_ltv = 0

    data['blended_arpu'] = blended_arpu
    data['blended_ltv'] = blended_ltv
    data['months_from_churn'] = months_to_churn
    data['margin_fy29'] = (data['net_income_fy29'] / data['revenue_fy29'] * 100) if data['revenue_fy29'] > 0 else 0

    wb.close()
    return data


def update_pitch_slides(model_path, slides_path, output_path):
    """Update pitch slides from financial model data."""
    data = extract_excel_data(model_path)
    prs = Presentation(slides_path)

    # Slide 2: Financial Projections — Revenue chart
    if len(prs.slides) > 1:
        slide = prs.slides[1]
        updates = [
            ("$49K", fmt_k(data['revenue_fy26'])),
            ("$390K", fmt_k(data['revenue_fy27'])),
            ("$1.45M", fmt_k(data['revenue_fy28'])),
            ("$4.04M", fmt_k(data['revenue_fy29'])),
            ("18 customers", f"{fmt_comma(data['total_customers_fy26'])} customers"),
            ("229 customers", f"{fmt_comma(data['total_customers_fy27'])} customers"),
            ("846 customers", f"{fmt_comma(data['total_customers_fy28'])} customers"),
            ("2,092 customers", f"{fmt_comma(data['total_customers_fy29'])} customers"),
        ]
        for old, new in updates:
            for shape in slide.shapes:
                replace_text_in_shape(shape, old, new)

        for shape in slide.shapes:
            if shape.has_chart:
                chart_data = CategoryChartData()
                chart_data.categories = ['FY 2026', 'FY 2027', 'FY 2028', 'FY 2029']
                chart_data.add_series('Revenue', [
                    data['revenue_fy26'], data['revenue_fy27'],
                    data['revenue_fy28'], data['revenue_fy29']
                ])
                shape.chart.replace_data(chart_data)

    # Slide 3: Customer Growth
    if len(prs.slides) > 2:
        slide = prs.slides[2]
        for shape in slide.shapes:
            replace_text_in_shape(shape, "0 → 2,099 Customers",
                                  f"0 → {int(data['total_customers_fy29'])} Customers")
            replace_text_in_shape(shape, f"FY29 Total: 1,550 businesses",
                                  f"FY29 Total: {fmt_comma(data['global_total_fy29'])} businesses")

        # Update stacked bar chart
        for shape in slide.shapes:
            if shape.has_chart:
                wb2 = load_workbook(model_path, data_only=True)
                ws_g = wb2["Geographic Split"]

                def geo_sum(er, gr, nr, col):
                    return (ws_g[f'{col}{er}'].value or 0) + \
                           (ws_g[f'{col}{gr}'].value or 0) + \
                           (ws_g[f'{col}{nr}'].value or 0)

                cols = ['B', 'C', 'D', 'E']
                chart_data = CategoryChartData()
                chart_data.categories = ['FY26', 'FY27', 'FY28', 'FY29']
                chart_data.add_series('Light', [geo_sum(6, 13, 20, c) for c in cols])
                chart_data.add_series('Growth', [geo_sum(7, 14, 21, c) for c in cols])
                chart_data.add_series('Pro', [geo_sum(8, 15, 22, c) for c in cols])
                chart_data.add_series('Enterprise', [geo_sum(9, 16, 23, c) for c in cols])
                shape.chart.replace_data(chart_data)
                wb2.close()

    # Slide 5: The Ask
    if len(prs.slides) > 4:
        slide = prs.slides[4]
        for shape in slide.shapes:
            replace_text_in_shape(shape, "Path to $4M ARR",
                                  f"Path to {fmt_k(data['revenue_fy29'])} ARR")

    # Slide 7: Unit Economics
    if len(prs.slides) > 6:
        slide = prs.slides[6]
        arpu = int(data['blended_arpu'])
        ltv = int(data['blended_ltv'])
        months = int(data['months_from_churn'])
        updates = [
            ("$166", f"${arpu:,}"), ("3%", fmt_pct(data['churn_rate'])),
            ("33 mo", f"{months} mo"), ("$5,478", f"${ltv:,}"),
        ]
        for old, new in updates:
            for shape in slide.shapes:
                replace_text_in_shape(shape, old, new)

    prs.save(output_path)
    return output_path
